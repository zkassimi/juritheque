"""
html_to_pptx.py
Convert presentation_juritheque.html → JuriTheque_Presentation.pptx
Each slide is screenshotted at 16:9 (1920×1080) and inserted as full-bleed image.
"""

import os, sys, io, time
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR  = Path(r"C:\Users\HP\Desktop\Legal-website\lexbase")
HTML_FILE = BASE_DIR / "presentation_juritheque.html"
OUT_FILE  = BASE_DIR / "JuriTheque_Presentation.pptx"
SHOTS_DIR = BASE_DIR / "_slide_shots"
SHOTS_DIR.mkdir(exist_ok=True)

# ── Dimensions 16:9 ──────────────────────────────────────────────────────────
W_PX, H_PX = 1920, 1080          # screenshot resolution
PPTX_W = Emu(9144000)            # 25.4 cm  (standard 16:9)
PPTX_H = Emu(5143500)            # 14.3 cm

def take_screenshots():
    screenshots = []
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": W_PX, "height": H_PX})

        url = HTML_FILE.as_uri()
        print(f"Loading: {url}")
        page.goto(url, wait_until="networkidle", timeout=30000)
        time.sleep(2)  # wait for fonts

        # Collect all .slide elements
        slides = page.query_selector_all(".slide")
        print(f"Found {len(slides)} slides")

        for i, slide_el in enumerate(slides):
            # Scroll into view & screenshot just that element
            slide_el.scroll_into_view_if_needed()
            time.sleep(0.3)

            png_path = SHOTS_DIR / f"slide_{i+1:02d}.png"
            slide_el.screenshot(path=str(png_path), type="png")
            screenshots.append(png_path)
            print(f"  ✓ slide {i+1:02d} → {png_path.name}")

        browser.close()
    return screenshots

def build_pptx(screenshots):
    prs = Presentation()
    prs.slide_width  = PPTX_W
    prs.slide_height = PPTX_H

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for i, img_path in enumerate(screenshots):
        slide = prs.slides.add_slide(blank_layout)
        # Insert image as full-bleed (cover entire slide)
        slide.shapes.add_picture(
            str(img_path),
            left=0, top=0,
            width=PPTX_W, height=PPTX_H
        )
        print(f"  ✓ added slide {i+1:02d} to PPTX")

    prs.save(str(OUT_FILE))
    print(f"\n✅ Saved: {OUT_FILE}")
    print(f"   {len(screenshots)} slides  |  {OUT_FILE.stat().st_size / 1024:.0f} KB")

if __name__ == "__main__":
    print("=" * 55)
    print("  HTML -> PPTX Converter - JuriTheque Presentation")
    print("=" * 55)

    if not HTML_FILE.exists():
        print(f"ERROR: HTML file not found: {HTML_FILE}")
        sys.exit(1)

    print("\n[1/2] Taking screenshots …")
    shots = take_screenshots()

    if not shots:
        print("ERROR: No slides found in the HTML file.")
        sys.exit(1)

    print(f"\n[2/2] Building PPTX …")
    build_pptx(shots)
    print("\nDone!")
