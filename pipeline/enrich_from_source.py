"""
enrich_from_source.py
─────────────────────
Pipeline en 2 étapes IA pour enrichir un guide depuis des PDF/PPTX.

Étape 1 — Filtre de pertinence (IA légère, ~200 tokens/doc)
  Pour chaque document du dossier, l'IA décide s'il est pertinent
  pour le guide cible. Les hors-sujet sont éliminés avant traitement.

Étape 2 — Génération des sections (IA principale)
  L'IA génère les sections FR + AR à partir des documents pertinents.

Usage :
  python pipeline/enrich_from_source.py --slug depenses-personnel-maroc --dir pipeline/sources/administratif/depenses-personnel/
  python pipeline/enrich_from_source.py --slug marches-publics-maroc --file cours.pdf
  python pipeline/enrich_from_source.py --slug urbanisme-maroc --dir pipeline/sources/urbanisme/ --dry-run

Options :
  --slug      Slug du guide cible (obligatoire)
  --file      Un ou plusieurs fichiers PDF/PPTX
  --dir       Dossier source (récursif — tous les sous-dossiers)
  --pages     Pages PDF ex: 10-45 (1 seul fichier uniquement)
  --lang      fr | ar | both  (défaut: both)
  --output    patch | print  (défaut: patch)
  --sections  Nombre de sections à générer (défaut: 3)
  --dry-run   Affiche le plan sans appeler l'IA de génération
  --no-filter Désactive le filtre de pertinence (traite tous les docs)
"""

import os, re, sys, json, argparse
from pathlib import Path
from datetime import datetime

import httpx
from dotenv import load_dotenv
from rich.console import Console
from rich.panel import Panel
from rich.syntax import Syntax
from rich.table import Table
from rich.progress import track

# ── Config ────────────────────────────────────────────────────────────────────
load_dotenv(Path(__file__).parent / ".env")

OPENROUTER_KEY = os.getenv("OPENROUTER_API_KEY") or os.getenv("VITE_OPENROUTER_KEY", "")
ANTHROPIC_KEY  = os.getenv("ANTHROPIC_API_KEY", "")
PROVIDER       = os.getenv("PROVIDER", "openrouter").lower()
MODEL_GENERATE = os.getenv("MODEL_GENERATE","google/gemini-2.5-pro")   # génération sections

ROOT        = Path(__file__).parent.parent
SEO_FILE    = ROOT / "src" / "data" / "seoIntentPages.js"
PATCHES_DIR = Path(__file__).parent / "patches"
PATCHES_DIR.mkdir(exist_ok=True)

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
MAX_CHARS_PER_DOC    = 6000   # chars envoyés par doc au filtre IA
MAX_CHARS_GENERATE   = 22000  # chars de source injectés à chaque appel
MAX_TOKENS_GENERATE  = 12000  # tokens réponse / appel (Gemini 2.5 Pro = modèle thinking ;
                              # le raisonnement consomme du budget → marge large pour 5-6 sections)
AUTO_MIN_SECTIONS    = 4      # mode auto (--sections 0) : borne basse
AUTO_MAX_SECTIONS    = 8      # mode auto : borne haute — l'IA choisit selon la richesse des sources

console = Console()

# ── Filtre PII minimal (noms, emails, téléphones) ────────────────────────────
_PII = [
    (re.compile(r'\b[\w.+-]+@[\w-]+\.[a-zA-Z]{2,}\b'), '[EMAIL]'),
    (re.compile(r'\b(?:\+212|00212|0)[567]\d[\s.\-]?\d{2}[\s.\-]?\d{2}[\s.\-]?\d{2}\b'), '[TEL]'),
    (re.compile(r'\b[A-Z]{1,2}\d{5,8}\b'), '[CIN]'),
    (re.compile(
        r'(?:^|\n)\s*(?:par|réalisé par|présenté par|préparé par|encadré par)\s*[:\-]?\s*'
        r'[A-ZÀ-Ü][A-Za-zÀ-ÿ\s]{2,40}(?=\n|$)',
        re.IGNORECASE | re.MULTILINE), ''),
]
_DROP_LINES = [
    re.compile(r'^\s*(?:université|faculté|école|fsjes|fsjp|encg)\b.*$', re.IGNORECASE),
    re.compile(r'^\s*(?:année universitaire|semestre [12]|filière\s*:).*$', re.IGNORECASE),
    re.compile(r'^\s*(?:page\s+\d+|\d+\s*/\s*\d+)\s*$', re.IGNORECASE),
    # Attribution d'auteur (PII)
    re.compile(r'(?:présenté|réalisé|préparé|encadré|élaboré|rédigé)\s+par\b', re.IGNORECASE),
    re.compile(r'^\s*par\s*[:\-]\s*[A-ZÀ-Ü]', re.IGNORECASE),
]

def sanitize(text: str) -> str:
    lines = []
    for line in text.splitlines():
        if any(p.search(line) for p in _DROP_LINES):
            continue
        for pat, repl in _PII:
            line = pat.sub(repl, line)
        if line.strip():
            lines.append(line)
    return "\n".join(lines)


def clean_extracted_text(text: str) -> str:
    """
    Nettoyage SÛR du texte extrait (PDF blocks / PPTX / DOCX).

    Principe : on ne FUSIONNE jamais de lignes entre elles (ça créait des
    phrases incohérentes). On se contente de retirer le bruit évident :
    - lignes purement numériques (axes de graphiques, pagination)
    - fragments isolés de 1-2 caractères
    - lignes vides multiples
    """
    out, blank = [], 0
    for line in text.splitlines():
        s = line.strip()
        if not s:
            blank += 1
            if blank <= 1:
                out.append("")
            continue
        blank = 0
        # Lignes uniquement chiffres / ponctuation / unités (graphiques, stats)
        if re.fullmatch(r'[\d\s,.\-•·%/()]+', s):
            continue
        # Fragments isolés de 1-2 caractères
        if len(s) <= 2:
            continue
        # Lignes décoratives à lettres espacées : "A C C È S  A  L A  F O N C T I O N"
        toks = s.split()
        if len(toks) >= 5 and sum(1 for t in toks if len(t) == 1) / len(toks) > 0.6:
            continue
        out.append(s)
    return "\n".join(out).strip()


# ── Extraction texte brut ─────────────────────────────────────────────────────

def extract_pdf_text(path: Path, pages_arg: str | None = None) -> str:
    try:
        import fitz, warnings
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            doc = fitz.open(str(path))
    except ImportError:
        console.print("[red]PyMuPDF non installé : pip install PyMuPDF[/]")
        sys.exit(1)
    except Exception as e:
        return ""

    total = len(doc)
    if pages_arg:
        try:
            parts = pages_arg.split("-")
            p0 = max(0, int(parts[0]) - 1)
            p1 = min(total, int(parts[1]) if len(parts) > 1 else total)
        except Exception:
            p0, p1 = 0, total
    else:
        p0, p1 = 0, total

    pages_text = []
    for i in range(p0, p1):
        try:
            page = doc[i]
            # Mode "blocks" : groupe le texte par bloc visuel (paragraphe / zone de
            # texte) → bien plus cohérent que "text" sur les PDF de présentation/slides.
            blocks = page.get_text("blocks")
            text_blocks = [b for b in blocks if b[6] == 0 and b[4].strip()]
            # Ordre de lecture : haut→bas (bande de 10pt), puis gauche→droite
            text_blocks.sort(key=lambda b: (round(b[1] / 10), round(b[0] / 10)))
            lines = []
            for b in text_blocks:
                # Joindre les retours-ligne internes d'un bloc (même paragraphe)
                btext = re.sub(r'\s*\n\s*', ' ', b[4].strip())
                btext = re.sub(r'[ \t]{2,}', ' ', btext)
                if btext:
                    lines.append(btext)
            if lines:
                pages_text.append("\n".join(lines))
        except Exception:
            pass
    doc.close()
    raw = sanitize("\n\n".join(pages_text))
    return clean_extracted_text(raw)


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        console.print("[red]python-pptx non installé : pip install python-pptx[/]")
        sys.exit(1)

    try:
        prs = Presentation(str(path))
    except Exception as e:
        console.print(f"   [yellow]⚠ PPTX ignoré ({path.name}) : {e}[/]")
        return ""

    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        # Titre via API officielle
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_texts.append(f"[{slide.shapes.title.text.strip()}]")
        except Exception:
            pass
        # Toutes les autres shapes
        for shape in slide.shapes:
            try:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                if shape.shape_type == 13:  # image
                    continue
                txt = shape.text.strip()
                # Ne pas dupliquer le titre
                if slide_texts and txt in slide_texts[0]:
                    continue
                slide_texts.append(txt)
            except Exception:
                continue
        if slide_texts:
            parts.append("\n".join(slide_texts))

    raw = sanitize("\n\n".join(parts))
    return clean_extracted_text(raw)


def extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        console.print("[red]python-docx non installé : pip install python-docx[/]")
        sys.exit(1)

    try:
        doc = Document(str(path))
    except Exception as e:
        console.print(f"   [yellow]⚠ DOCX ignoré ({path.name}) : {e}[/]")
        return ""

    parts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    # Tableaux (barèmes, grilles indiciaires, échelons...)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    raw = sanitize("\n".join(parts))
    return clean_extracted_text(raw)


def extract_text(path: Path, pages_arg: str | None = None) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path, pages_arg)
    elif suffix == ".pptx":
        return extract_pptx_text(path)
    elif suffix == ".docx":
        return extract_docx_text(path)
    elif suffix == ".ppt":
        console.print(f"   [yellow]⚠ .ppt ignoré ({path.name}) — convertir en .pptx[/]")
        return ""
    elif suffix == ".doc":
        console.print(f"   [yellow]⚠ .doc ignoré ({path.name}) — convertir en .docx[/]")
        return ""
    return ""


# ── Collecte des fichiers sources ─────────────────────────────────────────────

def collect_files(file_args: list, dir_arg: str | None) -> list:
    paths = []
    for f in (file_args or []):
        p = Path(f)
        if not p.exists():
            console.print(f"[red]Fichier introuvable : {p}[/]"); sys.exit(1)
        if p.suffix.lower() in SUPPORTED_EXTENSIONS:
            paths.append(p)

    if dir_arg:
        d = Path(dir_arg)
        if not d.is_dir():
            console.print(f"[red]Dossier introuvable : {d}[/]"); sys.exit(1)
        for p in sorted(d.rglob("*")):
            if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS:
                paths.append(p)

    if not paths:
        console.print("[red]Aucun fichier source trouvé.[/]"); sys.exit(1)

    # Dédoublonnage
    seen, unique = set(), []
    for p in paths:
        k = p.resolve()
        if k not in seen:
            seen.add(k); unique.append(p)
    return unique


# ── Lecture du guide ──────────────────────────────────────────────────────────

def read_guide(slug: str) -> dict:
    if not SEO_FILE.exists():
        console.print(f"[red]Introuvable : {SEO_FILE}[/]"); sys.exit(1)

    content = SEO_FILE.read_text(encoding="utf-8")
    m = re.search(rf"slug:\s*['\"]({re.escape(slug)})['\"]", content)
    if not m:
        console.print(f"[red]Guide introuvable : {slug}[/]")
        slugs = re.findall(r"slug:\s*['\"]([^'\"]+)['\"]", content)
        console.print("[yellow]" + "\n".join(f"  - {s}" for s in slugs)); sys.exit(1)

    # Borne le bloc au guide réel (du slug jusqu'au slug suivant) — sinon une
    # fenêtre fixe trop courte rate sections_ar/faq_ar sur les gros guides.
    start = max(0, m.start() - 200)
    next_m = re.search(r"slug:\s*['\"]", content[m.end():])
    end = m.end() + next_m.start() if next_m else len(content)
    block = content[start:end]

    def field(name):
        m2 = re.search(rf"{name}:\s*['\"]([^'\"]*)['\"]", block)
        return m2.group(1) if m2 else ""

    def lst(name):
        m2 = re.search(rf"{name}:\s*\[([^\]]*)\]", block, re.DOTALL)
        return re.findall(r"['\"]([^'\"]{4,})['\"]", m2.group(1)) if m2 else []

    existing_h2 = re.findall(r"h2:\s*['\"]([^'\"]+)['\"]", block)

    return {
        "slug":            slug,
        "h1":              field("h1"),
        "h1_ar":           field("h1_ar"),
        "intro":           field("intro"),
        "legalDomain":     field("legalDomain"),
        "category":        field("category"),
        "keywords":        lst("keywords"),
        "keywords_ar":     lst("keywords_ar"),
        "existing_h2":     existing_h2,
        "has_sections":    bool(re.search(r"\bsections:", block)),
        "has_sections_ar": bool(re.search(r"\bsections_ar:", block)),
        "has_faq_ar":      bool(re.search(r"\bfaq_ar:", block)),
    }


# ── Extraction JSON robuste (gère les fences markdown + trailing commas) ─────

def parse_ai_json(raw: str) -> dict:
    """
    Extrait un objet JSON de la réponse IA de façon robuste.

    1. Nettoie les fences markdown
    2. json.loads strict (préféré si le JSON est déjà valide)
    3. json-repair : répare guillemets non échappés, virgules traînantes,
       JSON tronqué — le problème classique du JSON généré par LLM
    """
    clean = re.sub(r"^```(?:json)?\s*", "", raw.strip())
    clean = re.sub(r"\s*```$", "", clean.strip())

    # 1. Essai strict (rapide, sans modification du contenu)
    try:
        return json.loads(clean)
    except json.JSONDecodeError:
        pass

    # 2. Réparation robuste (guillemets internes, virgules, troncatures)
    try:
        from json_repair import repair_json
        repaired = repair_json(clean, return_objects=True)
        if isinstance(repaired, dict) and repaired:
            return repaired
    except ImportError:
        console.print("[yellow]⚠ json-repair absent : pip install json-repair[/]")
    except Exception:
        pass

    # 3. Dernier recours : extraire entre premier { et dernier }
    start, end = clean.find('{'), clean.rfind('}')
    if start != -1 and end > start:
        candidate = re.sub(r',\s*([}\]])', r'\1', clean[start:end + 1])
        return json.loads(candidate)

    raise json.JSONDecodeError("Aucun objet JSON trouvé dans la réponse", raw, 0)


# ── Appel IA générique ────────────────────────────────────────────────────────

def call_ai(prompt: str, max_tokens: int, model: str | None = None) -> str:
    m = model or MODEL_GENERATE

    if not OPENROUTER_KEY:
        console.print("[red]❌ OPENROUTER_API_KEY manquant dans pipeline/.env[/]"); sys.exit(1)

    r = httpx.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {OPENROUTER_KEY}", "Content-Type": "application/json"},
        json={"model": m, "messages": [{"role": "user", "content": prompt}],
              "max_tokens": max_tokens, "temperature": 0.2,
              "response_format": {"type": "json_object"}},
        timeout=180,
    )
    if not r.is_success:
        console.print(f"[red]❌ OpenRouter {r.status_code} : {r.text[:300]}[/]")
        r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"].strip()


# ── Étape 1 : Filtre de pertinence par mots-clés (Python pur, 0 token) ────────

def keyword_filter(guide: dict, text: str) -> dict:
    """
    Score de pertinence basé sur les mots-clés du guide.
    Rapide, gratuit, fiable — aucun appel API.
    """
    text_lower = text.lower()

    # Termes du guide : mots-clés + mots significatifs du h1
    kw_terms  = [kw.lower() for kw in guide["keywords"][:12]]
    h1_words  = [w.lower() for w in re.split(r'\W+', guide["h1"]) if len(w) > 4]
    all_terms = list(dict.fromkeys(kw_terms + h1_words))  # dédoublonnage, ordre conservé

    hits   = [t for t in all_terms if t in text_lower]
    score  = len(hits) / max(len(all_terms), 1)
    conf   = min(round(score * 2, 2), 1.0)    # score × 2 → confidence (plafonné à 1)

    return {
        "relevant":   score >= 0.12,           # au moins ~1-2 termes sur 10-12
        "confidence": conf,
        "reason":     f"{len(hits)}/{len(all_terms)} termes — {', '.join(hits[:4])}",
        "topics":     ", ".join(hits[:6]),
    }


# ── Étape 2 : Génération des sections ────────────────────────────────────────

def build_source_context(relevant_docs: list) -> str:
    """Concatène les sources, budget RÉPARTI équitablement entre documents."""
    source_blocks = []
    per_doc = max(MAX_CHARS_GENERATE // max(len(relevant_docs), 1), 2500)
    for d in relevant_docs:
        chunk = d["text"][:per_doc]
        label = f"📄 {d['path'].name}"
        if d.get("topics"):
            label += f" — sujets : {d['topics']}"
        source_blocks.append(f"### {label}\n{chunk}")
    return "\n".join(source_blocks)


def plan_targets(guide: dict, lang: str, force: bool) -> list:
    """Clés à générer selon la langue et l'état du guide (faq_ar seulement si absente)."""
    targets = []
    if lang in ("fr", "both") and (force or not guide["has_sections"]):
        targets.append("sections")
    if lang in ("ar", "both") and (force or not guide["has_sections_ar"]):
        targets.append("sections_ar")
    if lang in ("ar", "both") and not guide["has_faq_ar"]:
        targets.append("faq_ar")
    return targets


def build_target_prompt(guide: dict, source_ctx: str, target: str,
                        n_sections: int, mirror_h2: list | None = None) -> str:
    """
    Prompt FOCALISÉ sur une seule cible (un appel = une langue) → tout le budget
    tokens pour cette sortie, donc zéro troncature.

    n_sections == 0 → mode AUTO : l'IA choisit le nombre (AUTO_MIN..AUTO_MAX).
    n_sections  > 0 → nombre EXACT imposé.
    mirror_h2       → (sections_ar) liste des h2 FR à refléter : même nombre,
                      mêmes sous-thèmes, même ordre, adaptés en arabe.
    """
    existing_note = ""
    if guide["existing_h2"]:
        existing_note = "H2 déjà présents ailleurs (ne pas dupliquer) :\n" + \
                        "\n".join(f"  - {h}" for h in guide["existing_h2"])

    header = f"""Tu es un juriste expert en droit public marocain. Tu rédiges le contenu de fond d'un guide juridique de juritheque.com.

## Guide cible
- Titre FR : {guide['h1']}
- Titre AR : {guide['h1_ar']}
- Domaine : {guide['legalDomain']}
- Intro : {guide['intro'][:300]}
- Mots-clés : {', '.join(guide['keywords'][:10])}
{existing_note}

## Documents sources
{source_ctx}
"""

    if target in ("sections", "sections_ar"):
        is_ar  = target == "sections_ar"
        langue = "ARABE (arabe standard moderne)" if is_ar else "FRANÇAIS"
        arabe_rule = (
            "Rédige un arabe juridique fluide et naturel. Si les sources contiennent de l'arabe "
            "mal encodé (mojibake), ignore-le et rédige toi-même un arabe correct à partir du contenu FRANÇAIS."
            if is_ar else
            "Style clair, précis et professionnel, accessible à un non-juriste."
        )

        # Consigne sur le NOMBRE de sections
        if mirror_h2:
            titres = "\n".join(f"  {i+1}. {h}" for i, h in enumerate(mirror_h2))
            count_rule = (
                f"Génère EXACTEMENT {len(mirror_h2)} sections en arabe, correspondant une à une "
                f"aux sections françaises ci-dessous (même sujet, même ordre) :\n{titres}\n"
                f"Pour chacune, rédige un \"h2\" arabe (adaptation naturelle, pas une traduction mot-à-mot) "
                f"et un contenu arabe complet et original."
            )
        elif n_sections and n_sections > 0:
            count_rule = f"Génère EXACTEMENT {n_sections} sections de fond en {langue}."
        else:
            count_rule = (
                f"Détermine TOI-MÊME le nombre de sections approprié, entre {AUTO_MIN_SECTIONS} et "
                f"{AUTO_MAX_SECTIONS}, selon la richesse RÉELLE des sources : une section par sous-thème "
                f"distinct effectivement couvert. N'invente JAMAIS de section creuse ou redondante pour "
                f"gonfler le nombre ; mieux vaut {AUTO_MIN_SECTIONS} sections solides que {AUTO_MAX_SECTIONS} diluées."
            )

        return header + f"""
## Tâche
{count_rule}

## Règles impératives
1. Chaque section porte sur un aspect DIFFÉRENT et COMPLÉMENTAIRE du sujet (ex. selon le sujet : accès/recrutement, avancement, positions, rémunération, sanctions, sortie de service, contrôle...).
2. Chaque section : "h2" précis et spécifique + "content" = tableau de 3 paragraphes d'environ 120 mots + "bullets" = 5 à 6 points concrets et informatifs.
3. Contenu juridique marocain RÉEL : cite précisément les lois/dahirs/décrets/articles présents dans les sources (n° et dates).
4. Va en PROFONDEUR : explique procédures, conditions, délais, montants, acteurs et organes. Évite les généralités vagues.
5. {arabe_rule}
6. Réponds en JSON valide UNIQUEMENT, sans markdown, sans commentaire.

## Format attendu
{{"{target}": [{{"h2": "...", "content": ["...", "...", "..."], "bullets": ["...", "...", "...", "...", "..."]}}]}}
"""

    # faq_ar
    return header + """
## Tâche
Rédige 5 questions/réponses fréquentes en ARABE (arabe standard moderne) pour ce guide.

## Règles impératives
1. Questions concrètes que se posent réellement citoyens et professionnels sur ce sujet au Maroc.
2. Réponses précises (2 à 4 phrases), fondées sur le droit marocain réel issu des sources.
3. Si les sources contiennent de l'arabe mal encodé, ignore-le et rédige un arabe correct.
4. Réponds en JSON valide UNIQUEMENT, sans markdown.

## Format attendu
{"faq_ar": [{"question": "...", "answer": "..."}]}
"""


# ── Convertir JSON → patch JS ─────────────────────────────────────────────────

def json_to_js(data: dict, slug: str) -> str:
    def js_str(s: str) -> str:
        return "'" + str(s).replace("\\", "\\\\").replace("'", "\\'") + "'"

    def js_list(lst: list) -> str:
        items = ",\n          ".join(js_str(i) for i in lst)
        return f"[\n          {items},\n        ]"

    def js_item(item: dict, pad: str) -> str:
        lines = [f"{pad}{{"]
        for k in ("h2", "question"):
            if k in item: lines.append(f"{pad}  {k}: {js_str(item[k])},")
        if "content" in item:
            c = item["content"]
            lines.append(f"{pad}  content: {js_list(c) if isinstance(c, list) else js_str(str(c))},")
        if "bullets" in item and item["bullets"]:
            lines.append(f"{pad}  bullets: {js_list(item['bullets'])},")
        if "answer" in item:
            lines.append(f"{pad}  answer: {js_str(item['answer'])},")
        lines.append(f"{pad}}}")
        return "\n".join(lines)

    date_str = datetime.now().strftime("%Y-%m-%d")
    out = [f"// ── PATCH {slug} — {date_str} ──────────────────────────",
           f"// Coller dans seoIntentPages.js dans l'objet '{slug}'", ""]

    for key, items in data.items():
        if not isinstance(items, list) or not items:
            continue
        out.append(f"    {key}: [")
        out.append(",\n".join(js_item(item, "      ") for item in items))
        out.append("    ],\n")

    return "\n".join(out)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--slug",      required=True)
    parser.add_argument("--file",      nargs="+", default=[])
    parser.add_argument("--dir",       default=None)
    parser.add_argument("--pages",     default=None)
    parser.add_argument("--lang",      default="both", choices=["fr", "ar", "both"])
    parser.add_argument("--output",    default="patch", choices=["patch", "print"])
    parser.add_argument("--sections",  default=0, type=int,
                        help="Nb de sections (0 = auto : l'IA décide entre 4 et 8 selon les sources)")
    parser.add_argument("--dry-run",   action="store_true")
    parser.add_argument("--no-filter", action="store_true",
                        help="Désactive le filtre de pertinence")
    parser.add_argument("--force",     action="store_true",
                        help="Régénère les sections même si déjà présentes")
    args = parser.parse_args()

    # ── 1. Lire le guide ────────────────────────────────────────────────────
    console.print(f"\n[bold cyan]📖 Guide : {args.slug}[/]")
    guide = read_guide(args.slug)
    console.print(f"   h1: [white]{guide['h1']}[/] | domaine: {guide['legalDomain']}")
    console.print(
        f"   sections: {'[green]✓[/]' if guide['has_sections'] else '[red]✗[/]'}  "
        f"sections_ar: {'[green]✓[/]' if guide['has_sections_ar'] else '[red]✗[/]'}  "
        f"faq_ar: {'[green]✓[/]' if guide['has_faq_ar'] else '[red]✗[/]'}"
    )

    # ── 2. Collecter les fichiers ────────────────────────────────────────────
    files = collect_files(args.file, args.dir)
    pages_arg = args.pages if len(files) == 1 else None
    console.print(f"\n[bold cyan]📁 {len(files)} fichier(s) trouvé(s)[/]")

    # ── 3. Extraction texte brut ─────────────────────────────────────────────
    console.print(f"\n[bold cyan]📄 Extraction du texte...[/]")
    docs = []
    for path in files:
        try:
            rel = path.relative_to(Path(args.dir)) if args.dir else path
        except ValueError:
            rel = path
        text = extract_text(path, pages_arg if path == files[0] else None)
        if len(text.strip()) < 100:
            console.print(f"   [dim]⊘ {rel} — vide ou illisible[/]")
            continue
        docs.append({"path": path, "rel": rel, "text": text})
        console.print(f"   ✓ {rel} — {len(text):,} chars")

    if not docs:
        console.print("[red]❌ Aucun texte extrait.[/]"); sys.exit(1)

    # ── 4. Filtre IA de pertinence ───────────────────────────────────────────
    if args.no_filter:
        relevant = [dict(d, topics="") for d in docs]
        console.print(f"\n[dim]Filtre désactivé — {len(relevant)} docs traités[/]")
    else:
        console.print(f"\n[bold cyan]🔍 Étape 1/2 — Filtre par mots-clés ({len(docs)} docs)...[/]")
        relevant = []
        skipped  = []

        table = Table(show_header=True, header_style="bold", box=None, padding=(0, 2))
        table.add_column("Fichier", style="white")
        table.add_column("Score", justify="center")
        table.add_column("Confiance", justify="center")
        table.add_column("Termes trouvés", style="dim")

        for d in docs:
            result = keyword_filter(guide, d["text"])
            icon = "[green]✓[/]" if result["relevant"] else "[red]✗[/]"
            conf = f"{result['confidence']:.0%}"
            table.add_row(str(d["rel"])[:55], icon, conf, result["reason"][:65])

            if result["relevant"]:
                relevant.append(dict(d, topics=result["topics"]))
            else:
                skipped.append(d["path"].name)

        console.print(table)
        console.print(f"\n   [green]{len(relevant)} pertinent(s)[/] | [dim]{len(skipped)} ignoré(s)[/]")
        if skipped:
            console.print(f"   [dim]Ignorés : {', '.join(skipped[:5])}{'...' if len(skipped)>5 else ''}[/]")

    if not relevant:
        console.print("[yellow]⚠ Aucun document pertinent trouvé pour ce guide.[/]")
        console.print("[dim]Astuce : utilise --no-filter pour forcer le traitement de tous les docs.[/]")
        sys.exit(0)

    # ── 5. Planifier les cibles de génération ────────────────────────────────
    targets = plan_targets(guide, args.lang, args.force)
    if not targets:
        console.print(f"[green]✓ Guide '{guide['slug']}' déjà complet pour lang={args.lang}.[/] "
                      f"[dim](--force pour régénérer)[/]")
        sys.exit(0)

    source_ctx   = build_source_context(relevant)
    total_source = sum(len(d["text"]) for d in relevant)
    nb_label = f"auto ({AUTO_MIN_SECTIONS}-{AUTO_MAX_SECTIONS})" if not args.sections else str(args.sections)
    console.print(f"\n[bold cyan]📝 Sources retenues : {len(relevant)} docs | {total_source:,} chars[/]")
    console.print(f"   Cibles : [white]{', '.join(targets)}[/] | sections : {nb_label}")

    if args.dry_run:
        demo = build_target_prompt(guide, source_ctx, targets[0], args.sections)
        console.print(Panel(
            demo[:5000] + ("\n[...tronqué]" if len(demo) > 5000 else ""),
            title=f"[yellow]DRY-RUN — Prompt « {targets[0]} »[/]", expand=False
        ))
        return

    # ── 6. Génération IA — UN APPEL PAR CIBLE (tout le budget par langue,
    #       donc plus de sections et zéro troncature) ─────────────────────────
    console.print(f"\n[bold cyan]🤖 Étape 2/2 — Génération ({PROVIDER} / {MODEL_GENERATE})...[/]")
    data = {}
    for target in targets:
        console.print(f"   [dim]→ {target}...[/]")
        # L'arabe reflète la structure FR déjà générée (même nombre, mêmes sous-thèmes)
        mirror = None
        if target == "sections_ar" and isinstance(data.get("sections"), list):
            mirror = [s.get("h2", "") for s in data["sections"] if s.get("h2")]
        prompt = build_target_prompt(guide, source_ctx, target, args.sections, mirror_h2=mirror)
        raw = call_ai(prompt, MAX_TOKENS_GENERATE)
        try:
            part = parse_ai_json(raw)
        except json.JSONDecodeError as e:
            console.print(f"[red]   ❌ JSON invalide ({target}) : {e}[/]")
            continue
        if isinstance(part, dict) and isinstance(part.get(target), list):
            data[target] = part[target]
        elif isinstance(part, list):
            data[target] = part
        elif isinstance(part, dict):
            for v in part.values():
                if isinstance(v, list):
                    data[target] = v
                    break

    if not data:
        console.print("[red]❌ Aucune section générée.[/]"); sys.exit(1)

    # ── 8. Output ────────────────────────────────────────────────────────────
    js_patch = json_to_js(data, args.slug)

    if args.output == "print":
        console.print(Panel(
            Syntax(js_patch, "javascript", theme="monokai"),
            title=f"[green]✅ Patch — {args.slug}[/]", expand=False
        ))
    else:
        patch_file = PATCHES_DIR / f"{args.slug}.patch.js"
        patch_file.write_text(js_patch, encoding="utf-8")
        console.print(f"\n[bold green]✅ Patch sauvegardé :[/] {patch_file}")
        console.print("   1. Coller dans seoIntentPages.js")
        console.print("   2. npm run build:full")

    # Aperçu des sections générées
    for key, items in data.items():
        if isinstance(items, list) and items:
            console.print(f"\n[dim]{key} ({len(items)}) :[/]")
            for item in items:
                t = item.get("h2") or item.get("question", "")
                console.print(f"   [dim]▸[/] {t[:70]}")


if __name__ == "__main__":
    main()
